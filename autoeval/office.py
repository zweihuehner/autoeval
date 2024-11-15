from pptx import Presentation
from pathlib import Path
import numpy as np
import pandas as pd

from .io import ComparisonInputs

class PowerPointCreator:
    
    def __init__(self, save_file: str, title: str, base_presentation: str):
        """
        Initializes the PowerPointCreator.

        Args:
            save_file (str): The path to save the presentation.
            title (str): The title of the presentation.
            base_presentation (str): The path to the base presentation.".
        """
        self.prs = Presentation(base_presentation)
        self.prs.slides[0].shapes.title.text = title
        self.save_file = save_file
        self.n_baseslides = len(self.prs.slides)

    def add_comparison_slide(self, file_scatter: Path, file_timeseries: Path, title: str = "Scatter and Timeseries", base_slide: str = "Comparison Base Slide"):
        """
        Adds a comparison slide with a scatter plot on the left, and a timeseries plot on the right to the PowerPoint presentation.

        Args:
            file_scatter (Path): The path to the image file for the scatter plot to be added to the slide.
            file_timeseries (Path): The path to the image file for the timeseries plot to be added to the slide.
            title (str): The title text for the slide. Defaults to 'Scatter and Timeseries'.
            base_slide (str, optional): The title of the base slide in the base presentation to use for layout. Defaults to 'Comparison Base Slide'.
        """
        base_slide = self._find_slide_by_title(title = base_slide)

        slide = self.prs.slides.add_slide(base_slide.slide_layout)
        slide.shapes.title.text = title

        left, top, width, height = self._get_shape_dimension(slide, idx = 22) # Idx 22 is the index of the scatter placeholder
        slide.shapes.add_picture(file_scatter.as_posix(), left = left, top = top, width=width, height=height)
        left, top, width, height = self._get_shape_dimension(slide, idx = 23) # Idx 23 is the index of the timeseries placeholder
        slide.shapes.add_picture(file_timeseries.as_posix(), left = left, top = top, width=width, height=height)

    def add_overview_slide(self, file_overview: Path, title: str = "Overview", base_slide: str = "Overview Base Slide"):
        """
        Adds an overview slide with a picture to the PowerPoint presentation.

        Args:
            file_overview (Path): The path to the image file for the overview plot to be added to the slide.
            title (str): The title text for the slide. Defaults to 'Overview'.
            base_slide (str, optional): The title of the base slide in the base presentation to use for layout. Defaults to 'Overview Base Slide'.
        """
        base_slide = self._find_slide_by_title(title = base_slide)

        slide = self.prs.slides.add_slide(base_slide.slide_layout)
        slide.shapes.title.text = title

        left, top, width, height = self._get_shape_dimension(slide, idx = 13) # Idx 13 is the index of the first picture placeholder
        slide.shapes.add_picture(file_overview.as_posix(), left = left, top = top, width=width, height=height)

    def add_inspection_slide(self, file_timeseries: Path, title: str = "Inspection", base_slide: str = "Overview Base Slide"):
        """
        Adds an inspection slide with a picture to the PowerPoint presentation.

        Args:
            file_timeseries (Path): The path to the image file for the inspection plot to be added to the slide.
            title (str): The title text for the slide. Defaults to 'Inspection'.
            base_slide (str, optional): The title of the base slide in the base presentation to use for layout. Defaults to 'Overview Base Slide'.
        """
        base_slide = self._find_slide_by_title(title = base_slide)
        
        slide = self.prs.slides.add_slide(base_slide.slide_layout)
        slide.shapes.title.text = title

        left, top, width, height = self._get_shape_dimension(slide, idx = 13) # Idx 13 is the index of the first picture placeholder
        slide.shapes.add_picture(file_timeseries.as_posix(), left = left, top = top, width=width, height=height)

    def create(self):
        """
        Creates the PowerPoint presentation by adding the slides created with the add functions to the presentation
        and saving the presentation to the specified file.

        The function first removes all base slides from the presentation except the first one. Then it adds the slides
        created with the add functions to the presentation. Finally, it saves the presentation to the specified file.
        """
        for _ in np.arange(1, self.n_baseslides):
            del self.prs.slides._sldIdLst[1]

        self.prs.save(self.save_file)

    def _get_shape_dimension(self, slide, idx: int):
        """
        Retrieves the dimension of a shape (placeholder) with a given idx from a slide and removes it.

        Args:
            slide (pptx.slide.Slide): The slide to get the shape (placeholder) from.
            idx (int): The idx of the shape (placeholder) to retrieve.

        Returns:
            tuple: A tuple containing the left, top, width, and height of the shape.
        """
        for shape in slide.shapes:
            if shape.shape_type == 14:  # Shape type 14 corresponds to placeholders
                if shape.placeholder_format.idx == idx:
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    slide.shapes._spTree.remove(shape._element)
                    return left, top, width, height

    def _find_slide_by_title(self, title):
        """
        Finds a slide by its title.

        Args:
            title (str): The title of the slide to find.

        Returns:
            Slide: The slide with the given title.

        Raises:
            ValueError: If no slide with the given title is found.
        """
        for slide in self.prs.slides:
            if slide.shapes.title and slide.shapes.title.text == title:
                return slide
        raise ValueError(f"No slide found with the title '{title}'")


def create_excel_table_from_df(data: ComparisonInputs,  folder_out: str, color_scale: str | list[str] | None = None):    
    """
    Creates an Excel table from a pandas DataFrame.

    Args:
        data (ComparisonInputs): The DataFrame to write to the Excel file.
        folder_out (str): The folder where the Excel file is saved.
        file_name (str, optional): The name of the Excel file. Defaults to "evaluation_table".
        color_scale (str | list[str] | None, optional): The color scale to use for conditional formatting. Defaults to None.

    Returns:
        None
    """
    color_scale_sets = {
            'green_to_white': {
                'type': '2_color_scale',
                'min_type': 'percent',
                'min_value': 0,
                'min_color': '#FFFFFF', 
                'max_type': 'percent',
                'max_value': 100,
                'max_color': '#43B75F'
                },
            'white_to_green': {
                'type': '2_color_scale',
                'min_type': 'percent',
                'min_value': 0,
                'min_color': '#43B75F', 
                'max_type': 'percent',
                'max_value': 100,
                'max_color': '#FFFFFF'
                },
            }

    df = data.skill

    allowed_color_scales = list(color_scale_sets.keys())
    if isinstance(color_scale, str):
        if color_scale not in allowed_color_scales:
            raise ValueError(f"color_scale_type must be one of {allowed_color_scales}")
    elif isinstance(color_scale, list):
        for color_scale_i in color_scale:
            if color_scale_i not in allowed_color_scales:
                raise ValueError(f"color_scale_type must be one of {allowed_color_scales}")
    elif color_scale is None:
        color_scale = ["white_to_green", "white_to_green", "white_to_green", "white_to_green", "green_to_white", "white_to_green", "green_to_white"]
    
    file_out = Path(folder_out) / f"evaluation_table_{data.quantity.replace(' ', '')}.xlsx"
    
    with pd.ExcelWriter(file_out, engine='xlsxwriter') as writer:
        df.to_excel(writer, sheet_name='Sheet1', index=True)

        workbook = writer.book
        worksheet = writer.sheets['Sheet1']
        num_rows, num_cols = df.shape

        header_format = workbook.add_format({
            'bg_color': '#5B6277',
            'font_color': '#FFFFFF',
            'bold': True,
            'align': 'center',
            'border': 1,
            'border_color': '#FFFFFF',
        })
        centered_format = workbook.add_format({
            'align': 'center', 
            'valign': 'vcenter',
            'border': 1
        })
        left_format = workbook.add_format({
            'align': 'left',  
            'valign': 'vcenter',
            'border': 1
        })

        worksheet.write(0, 0, df.index.name if df.index.name else 'Index', header_format)
        for col_num, value in enumerate(df.columns.values):
            worksheet.write(0, col_num+1, value, header_format)
        
        start_row = 2

        for row_num in range(len(df)):
            worksheet.write(row_num + 1, 0, df.index[row_num], left_format)
            for col_num, column_name in enumerate(df.columns):
                worksheet.write(row_num + start_row-1, col_num + 1, df[column_name].iloc[row_num], centered_format)
        
        if isinstance(color_scale, str):
            start_col = 1  
            data_range = f"{chr(65 + start_col)}{start_row}:{chr(65 + start_col + num_cols - 1)}{start_row + num_rows - 1}"
            worksheet.conditional_format(data_range, color_scale_sets[color_scale])
        else:

            assert len(color_scale) == num_cols, f"color_scale must be of length {num_cols}"
            for i, _ in enumerate(df.columns):
                col_letter = chr(65 + i + 1)
                data_range = f"{col_letter}{start_row}:{col_letter}{start_row + len(df) - 1}"
                worksheet.conditional_format(data_range, color_scale_sets[color_scale[i]])

        for col_num, column_name in enumerate(df.columns):
            max_length = max(df[column_name].astype(str).map(len).max(), len(column_name))
            worksheet.set_column(col_num, col_num, max_length + 2)

        max_index_length = len(df.index.name) if df.index.name else 0
        max_index_length = max(max_index_length, df.index.astype(str).map(len).max()) 
        worksheet.set_column(0, 0, max_index_length + 2) 
